""" This module helps creating specific type of slides using our template powerpoint using python-pptx """
import os
import sys
import os_util

from PIL import Image
from pptx import Presentation

# CONSTANTS
# HEIGHT = 9
# WIDTH = 16
# LEFTMOST = Inches(0)
# TOPMOST = Inches(0)
# HEIGHT_IN = Inches(HEIGHT)
# WIDTH_IN = Inches(WIDTH)

# One inch equates to 914400 EMUs
# INCHES_TO_EMU = 914400
# One centimeter is 360000 EMUs
# CMS_TO_EMU = 360000

# Location of powerpoint template
POWERPOINT_TEMPLATE_FILE = 'data/powerpoint/template.pptx'

# Layouts index in template
LAYOUT_TITLE_SLIDE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_SECTION_HEADER = 2
LAYOUT_TWO_CONTENT = 3
LAYOUT_TWO_TITLE_AND_CONTENT = 4
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6
LAYOUT_CONTENT_CAPTION = 7
LAYOUT_PICTURE_CAPTION = 8
LAYOUT_FULL_PICTURE = 11
LAYOUT_TITLE_AND_PICTURE = 12
LAYOUT_LARGE_QUOTE = 13
LAYOUT_TWO_TITLE_AND_IMAGE = 14
LAYOUT_THREE_TITLE_AND_IMAGE = 15

_PROHIBITED_IMAGES_DIR = "./data/images/prohibited/"
PROHIBITED_IMAGES = list([os_util.open_image(_PROHIBITED_IMAGES_DIR + url) for url in os.listdir(_PROHIBITED_IMAGES_DIR)])


# = HELPERS =

# VALIDITY CHECKING

def _is_image(content):
    lower_url = content.lower()
    return ".jpg" in lower_url or ".gif" in lower_url or ".png" in lower_url or ".jpeg" in lower_url


def _is_valid_image(image_url):
    try:
        im = os_util.open_image(image_url)
        if im in PROHIBITED_IMAGES:
            print(image_url, " IS DENIED")
            return False
        print(image_url, "IS ALLOWED")
    except OSError:
        return False

    return True


def _is_valid_content(content):
    if _is_image(content):
        return _is_valid_image(content)
    return bool(content)


# CREATION
def _create_slide(prs, slide_type):
    """ Creates a new slide in the given presentation using the slide_type template """
    return prs.slides.add_slide(prs.slide_layouts[slide_type])


def _add_title(slide, title):
    """ Adds the given title to the slide if the title is present"""
    if title:
        title_object = slide.shapes.title
        title_object.text = title
        return True


def _add_text(slide, placeholder_id, text):
    if text:
        placeholder = slide.placeholders[placeholder_id]
        placeholder.text = text
        return True


def _add_image(slide, placeholder_id, image_url, original_image_size=True):
    if not os.path.isfile(image_url):
        return None
    placeholder = slide.placeholders[placeholder_id]
    if original_image_size:
        # Calculate the image size of the image
        im = os_util.open_image(image_url)
        width, height = im.size

        # Make sure the placeholder doesn't zoom in
        placeholder.height = height
        placeholder.width = width

        # Insert the picture
        placeholder = placeholder.insert_picture(image_url)

        # Calculate ratios and compare
        image_ratio = width / height
        placeholder_ratio = placeholder.width / placeholder.height
        ratio_difference = placeholder_ratio - image_ratio

        # Placeholder width too wide:
        if ratio_difference > 0:
            difference_on_each_side = ratio_difference / 2
            placeholder.crop_left = -difference_on_each_side
            placeholder.crop_right = -difference_on_each_side
        # Placeholder height too high
        else:
            difference_on_each_side = -ratio_difference / 2
            placeholder.crop_bottom = -difference_on_each_side
            placeholder.crop_top = -difference_on_each_side

        return True
    else:
        try:
            placeholder.insert_picture(image_url)
        except OSError:
            print("Unexpected error inserting image:", image_url, ":", sys.exc_info()[0])
            return None


def _add_image_or_text(slide, placeholder_id, image_url_or_text, original_image_size):
    if _is_image(image_url_or_text):
        return _add_image(slide, placeholder_id, image_url_or_text, original_image_size)
    else:
        return _add_text(slide, placeholder_id, image_url_or_text)


def _print_all_placeholders(slide):
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


# FORMAT GENERATORS
# These are functions that get some inputs (texts, images...)
# and create layouted slides with these inputs

def create_new_powerpoint():
    return Presentation(POWERPOINT_TEMPLATE_FILE)


def create_title_slide(prs, title, subtitle):
    slide = _create_slide(prs, LAYOUT_TITLE_SLIDE)
    _add_title(slide, title)
    _add_text(slide, 1, subtitle)
    return slide


def create_large_quote_slide(prs, text):
    if bool(text):
        slide = _create_slide(prs, LAYOUT_LARGE_QUOTE)
        _add_text(slide, 1, text)
        return slide


def create_image_slide(prs, title=None, image_url=None, original_image_size=True):
    """ Creates a slide with an image covering the whole slide"""
    # TODO debug this: the image can not be set!
    return _create_single_image_slide(prs, title, image_url, LAYOUT_TITLE_AND_PICTURE, original_image_size)


def create_full_image_slide(prs, title=None, image_url=None, original_image_size=True):
    """ Creates a slide with an image covering the whole slide"""
    return _create_single_image_slide(prs, title, image_url, LAYOUT_FULL_PICTURE, original_image_size)


def create_two_column_images_slide(prs, title=None, caption_1=None, image_or_text_1=None, caption_2=None,
                                   image_or_text_2=None, original_image_size=True):
    if _is_valid_content(image_or_text_1) and _is_valid_content(image_or_text_2):
        slide = _create_slide(prs, LAYOUT_TWO_TITLE_AND_IMAGE)
        _add_title(slide, title)
        _add_text(slide, 1, caption_1)
        _add_image_or_text(slide, 13, image_or_text_1, original_image_size)
        _add_text(slide, 3, caption_2)
        _(slide, 14, image_or_text_2, original_image_size)
        return slide


def create_three_column_images_slide(prs, title=None, caption_1=None, image_or_text_1=None, caption_2=None,
                                     image_or_text_2=None, caption_3=None, image_or_text_3=None,
                                     original_image_size=True):
    if _is_valid_content(image_or_text_1) and _is_valid_content(image_or_text_2) and _is_valid_content(image_or_text_3):
        slide = _create_slide(prs, LAYOUT_THREE_TITLE_AND_IMAGE)
        _add_title(slide, title)
        _add_text(slide, 1, caption_1)
        _add_image_or_text(slide, 13, image_or_text_1, original_image_size)
        _add_text(slide, 3, caption_2)
        _add_image_or_text(slide, 14, image_or_text_2, original_image_size)
        _add_text(slide, 15, caption_3)
        _add_image_or_text(slide, 16, image_or_text_3, original_image_size)
        return slide


# def create_two_column_images_slide_text_second(prs, title=None, caption_1=None, image_1=None, caption_2=None,
#                                                quote=None,
#                                                original_image_size=True):
#     if bool(image_1):
#         slide = _create_slide(prs, LAYOUT_TWO_TITLE_AND_IMAGE)
#         _add_title(slide, title)
#         _add_text(slide, 1, caption_1)
#         _add_image_or_text(slide, 13, image_1, original_image_size)
#         _add_text(slide, 3, caption_2)
#         _add_image_or_text(slide, 14, quote)
#         return slide


def _create_single_image_slide(prs, title, image_url, slide_template_idx, fit_image):
    if _is_valid_content(image_url):
        slide = _create_slide(prs, slide_template_idx)
        _add_title(slide, title)
        _add_image_or_text(slide, 1, image_url, fit_image)
        return slide


# GENERATORS: Same as the template fillers above, but using generation functions

def generate_full_image_slide(title_generator, image_generator, original_image_size=True):
    return generate_slide(create_full_image_slide, (title_generator, image_generator, lambda x: original_image_size))


def generate_image_slide(title_generator, image_generator, original_image_size=True):
    return generate_slide(create_image_slide, (title_generator, image_generator, lambda x: original_image_size))


def generate_title_slide(title_generator, subtitle_generator):
    return generate_slide(create_title_slide, (title_generator, subtitle_generator))


def generate_large_quote_slide(text_generator):
    return generate_slide(create_large_quote_slide, (text_generator,))


def generate_two_column_images_slide(title_generator, caption_1_generator, image_1_generator, caption_2_generator,
                                     image_2_generator):
    return generate_slide(create_two_column_images_slide, (title_generator, caption_1_generator,
                                                           image_1_generator, caption_2_generator,
                                                           image_2_generator))


def generate_two_column_images_slide_tuple_caption(title_generator, captions_generator, image_1_generator,
                                                   image_2_generator):
    def generate(presentation_context, used):
        generated_tuple = captions_generator(presentation_context)
        generated = title_generator(presentation_context), generated_tuple[0], image_1_generator(
            presentation_context), generated_tuple[1], image_2_generator(presentation_context)

        if _is_different_enough(generated, used):
            return create_two_column_images_slide(get_presentation(presentation_context), *generated), generated

    return generate


def generate_three_column_images_slide(title_generator, caption_1_generator, image_1_generator, caption_2_generator,
                                       image_2_generator, caption_3_generator, image_3_generator):
    return generate_slide(create_three_column_images_slide, (title_generator, caption_1_generator,
                                                             image_1_generator, caption_2_generator,
                                                             image_2_generator, caption_3_generator,
                                                             image_3_generator))


# HELPERS

def generate_slide(slide_template, generators):
    def generate(presentation_context, used):
        generated = [content_generator(presentation_context) for content_generator in generators]
        if _is_different_enough(generated, used):
            return slide_template(get_presentation(presentation_context), *generated), generated

    return generate


def get_presentation(presentation_context):
    return presentation_context["presentation"]


def _is_different_enough(generated, used):
    (used_elements, allowed_repeated_elements) = used
    intersection = set(generated) & used_elements
    return allowed_repeated_elements >= len(intersection)
