import os
import sys
import traceback
from functools import lru_cache

from lxml.etree import XMLSyntaxError
from pptx import Presentation


from talkgenerator.util import os_util
# Location of powerpoint template
_POWERPOINT_TEMPLATE_FILE = '../../data/powerpoint/template.pptx'


@lru_cache(maxsize=1)
def get_powerpoint_template_file():
    return os_util.to_actual_file(_POWERPOINT_TEMPLATE_FILE, __file__)


# Layouts index in template
LAYOUT_TITLE_SLIDE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_SECTION_HEADER = 2
LAYOUT_TWO_CONTENT = 3
LAYOUT_TWO_TITLE_AND_CONTENT = 4
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6
LAYOUT_CONTENT_CAPTION = 7
LAYOUT_PICTURE_CAPTION = 8
LAYOUT_FULL_PICTURE = 11
LAYOUT_TITLE_AND_PICTURE = 12
LAYOUT_LARGE_QUOTE = 13
LAYOUT_TWO_TITLE_AND_IMAGE = 14
LAYOUT_THREE_TITLE_AND_IMAGE = 15
LAYOUT_TITLE_AND_CHART = 16


# = HELPERS =

# VALIDITY CHECKING


def _is_valid_content(content):
    if not bool(content):
        return False
    if os_util.is_image(content):
        return os_util.is_valid_image(content)
    return True


# CREATION
def _create_slide(prs, slide_type):
    """ Creates a new slide in the given presentation using the slide_type template """
    return prs.slides.add_slide(prs.slide_layouts[slide_type])


def _add_title(slide, title):
    """ Adds the given title to the slide if the title is present"""
    if title:
        title_object = slide.shapes.title
        title_object.text = title
        return True


def _add_text(slide, placeholder_id, text):
    if text:
        placeholder = slide.placeholders[placeholder_id]
        placeholder.text = str(text)
        return True


def _add_image(slide, placeholder_id, image_url, original_image_size=True):
    if not os.path.isfile(image_url):
        return None

    image_url = os_util.to_actual_file(image_url, __file__)

    placeholder = slide.placeholders[placeholder_id]
    if original_image_size:
        # Calculate the image size of the image
        try:
            im = os_util.open_image(image_url)
            width, height = im.size

            # Make sure the placeholder doesn't zoom in
            placeholder.height = height
            placeholder.width = width

            # Insert the picture
            try:
                placeholder = placeholder.insert_picture(image_url)
            except (ValueError, XMLSyntaxError) as e:
                traceback.print_exc(file=sys.stdout)
                print('_add_image error: {}'.format(e))
                return None

            # Calculate ratios and compare
            image_ratio = width / height
            placeholder_ratio = placeholder.width / placeholder.height
            ratio_difference = placeholder_ratio - image_ratio

            # Placeholder width too wide:
            if ratio_difference > 0:
                difference_on_each_side = ratio_difference / 2
                placeholder.crop_left = -difference_on_each_side
                placeholder.crop_right = -difference_on_each_side
            # Placeholder height too high
            else:
                difference_on_each_side = -ratio_difference / 2
                placeholder.crop_bottom = -difference_on_each_side
                placeholder.crop_top = -difference_on_each_side

            return placeholder
        except FileNotFoundError as fnfe:
            traceback.print_exc(file=sys.stdout)
            print('_add_image file not found: {}'.format(fnfe))
            return None
    else:
        try:
            return placeholder.insert_picture(image_url)
        except OSError or ValueError:
            traceback.print_exc(file=sys.stdout)
            print("Unexpected error inserting image:", image_url, ":", sys.exc_info()[0])
            return None


def _add_chart(slide, placeholder_id, chart_type, chart_data):
    placeholder = slide.placeholders[placeholder_id]
    return placeholder.insert_chart(chart_type, chart_data)


def _add_image_or_text(slide, placeholder_id, image_url_or_text, original_image_size):
    if os_util.is_image(image_url_or_text):
        return _add_image(slide, placeholder_id, image_url_or_text, original_image_size)
    else:
        return _add_text(slide, placeholder_id, image_url_or_text)


def _print_all_placeholders(slide):
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


# FORMAT GENERATORS
# These are functions that get some inputs (texts, images...)
# and create layouted slide with these inputs

def create_new_powerpoint():
    return Presentation(get_powerpoint_template_file())


def create_title_slide(prs, title, subtitle):
    slide = _create_slide(prs, LAYOUT_TITLE_SLIDE)
    _add_title(slide, title)
    _add_text(slide, 1, subtitle)
    return slide


def create_large_quote_slide(prs, title, text, background_image=None):
    if bool(text):
        slide = _create_slide(prs, LAYOUT_LARGE_QUOTE)
        if title:
            _add_title(slide, title)
        _add_text(slide, 1, text)
        if background_image:
            _add_image(slide, 10, background_image, False)

        # Add black transparent image for making other image behind it transparent (missing feature in python-pptx)
        _add_image(slide, 11, "../../data/images/black-transparent.png", False)

        return slide


def create_image_slide(prs, title=None, image_url=None, original_image_size=True):
    """ Creates a slide with an image covering the whole slide"""
    # TODO debug this: the image can not be set!
    return _create_single_image_slide(prs, title, image_url, LAYOUT_TITLE_AND_PICTURE, original_image_size)


def create_full_image_slide(prs, title=None, image_url=None, original_image_size=True):
    """ Creates a slide with an image covering the whole slide"""
    return _create_single_image_slide(prs, title, image_url, LAYOUT_FULL_PICTURE, original_image_size)


def create_two_column_images_slide(prs, title=None, caption_1=None, image_or_text_1=None, caption_2=None,
                                   image_or_text_2=None, original_image_size=True):
    if _is_valid_content(image_or_text_1) and _is_valid_content(image_or_text_2):
        slide = _create_slide(prs, LAYOUT_TWO_TITLE_AND_IMAGE)
        _add_title(slide, title)
        _add_text(slide, 1, caption_1)
        _add_image_or_text(slide, 13, image_or_text_1, original_image_size)
        _add_text(slide, 3, caption_2)
        _add_image_or_text(slide, 14, image_or_text_2, original_image_size)
        return slide


def create_three_column_images_slide(prs, title=None, caption_1=None, image_or_text_1=None, caption_2=None,
                                     image_or_text_2=None, caption_3=None, image_or_text_3=None,
                                     original_image_size=True):
    if _is_valid_content(image_or_text_1) and _is_valid_content(image_or_text_2) and _is_valid_content(image_or_text_3):
        slide = _create_slide(prs, LAYOUT_THREE_TITLE_AND_IMAGE)
        _add_title(slide, title)
        _add_text(slide, 1, caption_1)
        _add_image_or_text(slide, 13, image_or_text_1, original_image_size)
        _add_text(slide, 3, caption_2)
        _add_image_or_text(slide, 14, image_or_text_2, original_image_size)
        _add_text(slide, 15, caption_3)
        _add_image_or_text(slide, 16, image_or_text_3, original_image_size)
        return slide


# def create_two_column_images_slide_text_second(prs, title=None, caption_1=None, image_1=None, caption_2=None,
#                                                quote=None,
#                                                original_image_size=True):
#     if bool(image_1):
#         slide = _create_slide(prs, LAYOUT_TWO_TITLE_AND_IMAGE)
#         _add_title(slide, title)
#         _add_text(slide, 1, caption_1)
#         _add_image_or_text(slide, 13, image_1, original_image_size)
#         _add_text(slide, 3, caption_2)
#         _add_image_or_text(slide, 14, quote)
#         return slide


def _create_single_image_slide(prs, title, image_url, slide_template_idx, fit_image):
    if _is_valid_content(image_url):
        slide = _create_slide(prs, slide_template_idx)
        _add_title(slide, title)
        _add_image_or_text(slide, 1, image_url, fit_image)
        return slide


def create_chart_slide(prs, title, chart_type, chart_data, chart_modifier=None):
    slide = _create_slide(prs, LAYOUT_TITLE_AND_CHART)
    _add_title(slide, title)
    chart = _add_chart(slide, 10, chart_type, chart_data).chart
    if chart_modifier:
        chart_modifier(chart, chart_data)
    return slide
